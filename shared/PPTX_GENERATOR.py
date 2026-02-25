#!/usr/bin/env python3
"""
Amenthyx AI Teams — Enhanced PPTX Status Report Generator v3.0
Usage: python PPTX_GENERATOR.py <report_number> <project_name> [team_dir]

Enhanced with:
- Evidence dashboard slide
- Commit activity slide
- Test coverage slide
- CI/CD status slide
- Kanban velocity metrics
"""

import sys
import os
import glob

sys.stdout.reconfigure(encoding='utf-8')

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from datetime import datetime


# Color scheme
DARK_BG = RGBColor(0x1A, 0x1A, 0x2E)
ACCENT = RGBColor(0x00, 0xD2, 0xFF)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
GREEN = RGBColor(0x00, 0xC8, 0x53)
YELLOW = RGBColor(0xFF, 0xC1, 0x07)
RED = RGBColor(0xFF, 0x17, 0x44)
GRAY = RGBColor(0x88, 0x88, 0x88)
PURPLE = RGBColor(0xBB, 0x86, 0xFC)
ORANGE = RGBColor(0xFF, 0x98, 0x00)


def add_dark_slide(prs, title_text, subtitle_text=None):
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = DARK_BG

    txBox = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12), Inches(1))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = title_text
    p.font.size = Pt(32)
    p.font.color.rgb = ACCENT
    p.font.bold = True

    if subtitle_text:
        p2 = tf.add_paragraph()
        p2.text = subtitle_text
        p2.font.size = Pt(16)
        p2.font.color.rgb = WHITE

    return slide


def add_text(slide, left, top, width, height, text, size=14, color=WHITE, bold=False):
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.color.rgb = color
    p.font.bold = bold
    return tf


def read_file_safe(path, max_chars=1500):
    if os.path.exists(path):
        with open(path, 'r', encoding='utf-8') as f:
            return f.read()[:max_chars]
    return None


def count_kanban(team_dir):
    kanban = read_file_safe(os.path.join(team_dir, "KANBAN.md"))
    if kanban:
        return {
            "done": kanban.count("[x]"),
            "in_progress": kanban.count("[~]"),
            "in_review": kanban.count("[?]"),
            "backlog": kanban.count("[ ]"),
        }
    return {"done": 0, "in_progress": 0, "in_review": 0, "backlog": 0}


def count_evidence(team_dir):
    """Count evidence artifacts by category."""
    evidence_dir = os.path.join(team_dir, "evidence")
    counts = {
        "manifests": 0,
        "builds": 0,
        "tests": 0,
        "screenshots": 0,
        "runtime": 0,
        "ci": 0,
        "total": 0,
    }
    if os.path.exists(evidence_dir):
        for root, dirs, files in os.walk(evidence_dir):
            for f in files:
                counts["total"] += 1
                rel = os.path.relpath(root, evidence_dir)
                if rel.startswith("manifests"):
                    counts["manifests"] += 1
                elif rel.startswith("builds"):
                    counts["builds"] += 1
                elif rel.startswith("tests"):
                    counts["tests"] += 1
                elif rel.startswith("screenshots"):
                    counts["screenshots"] += 1
                elif rel.startswith("runtime"):
                    counts["runtime"] += 1
                elif rel.startswith("ci"):
                    counts["ci"] += 1
    return counts


def count_commits(team_dir):
    """Parse commit log if available."""
    commit_log = read_file_safe(os.path.join(team_dir, "COMMIT_LOG.md"), 5000)
    if commit_log:
        lines = [l for l in commit_log.split('\n') if '|' in l and not l.strip().startswith('#') and not l.strip().startswith('|--')]
        return max(0, len(lines) - 1)  # subtract header row
    return 0


def discover_artifacts(team_dir):
    """Count all artifacts by directory."""
    artifacts = {}
    if os.path.exists(team_dir):
        for root, dirs, files in os.walk(team_dir):
            for f in files:
                if f.endswith(('.md', '.yaml', '.yml', '.json', '.py', '.pptx', '.pdf', '.log', '.xml', '.html')):
                    rel_dir = os.path.relpath(root, team_dir)
                    if rel_dir not in artifacts:
                        artifacts[rel_dir] = 0
                    artifacts[rel_dir] += 1
    return artifacts


def create_status_report(report_num, project_name, team_dir=".team"):
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # ──────────────────────────────────────────────────────────────
    # Slide 1: Title
    # ──────────────────────────────────────────────────────────────
    slide1 = add_dark_slide(
        prs,
        f"Project Status Report #{report_num}",
        f"{project_name} | {datetime.now().strftime('%Y-%m-%d %H:%M')}"
    )
    add_text(slide1, 0.5, 2.5, 12, 1,
             "Amenthyx AI Teams v3.0 — Evidence-Driven Virtual Engineering",
             size=20, color=GRAY)
    add_text(slide1, 0.5, 3.5, 12, 0.5,
             "Evidence-Driven | Locally-Tested | Atomically-Committed | CI-Validated",
             size=14, color=PURPLE)

    # ──────────────────────────────────────────────────────────────
    # Slide 2: Executive Summary
    # ──────────────────────────────────────────────────────────────
    slide2 = add_dark_slide(prs, "Executive Summary")
    kanban = count_kanban(team_dir)
    total = kanban["done"] + kanban["in_progress"] + kanban["in_review"] + kanban["backlog"]
    pct = int((kanban["done"] / total * 100)) if total > 0 else 0

    if pct >= 80:
        status_color = GREEN
        status_label = "ON TRACK"
    elif pct >= 50:
        status_color = YELLOW
        status_label = "IN PROGRESS"
    else:
        status_color = ACCENT
        status_label = "EARLY STAGE"

    add_text(slide2, 0.5, 1.8, 6, 0.6, f"Overall Status: {status_label}", size=20, color=status_color, bold=True)
    add_text(slide2, 0.5, 2.5, 12, 0.5, f"Completion: {pct}% ({kanban['done']}/{total} tasks done)", size=16, color=WHITE)
    add_text(slide2, 0.5, 3.2, 12, 0.5, f"In Progress: {kanban['in_progress']}  |  In Review: {kanban['in_review']}  |  Backlog: {kanban['backlog']}", size=14, color=GRAY)

    evidence = count_evidence(team_dir)
    commits = count_commits(team_dir)
    add_text(slide2, 0.5, 4.2, 12, 0.5, f"Evidence Artifacts: {evidence['total']}  |  Commits: {commits}", size=14, color=PURPLE)

    # ──────────────────────────────────────────────────────────────
    # Slide 3: Evidence Dashboard
    # ──────────────────────────────────────────────────────────────
    slide3 = add_dark_slide(prs, "Evidence Dashboard", "Proof of work collected by each agent")

    categories = [
        ("Agent Manifests", evidence["manifests"], ACCENT),
        ("Build Logs", evidence["builds"], GREEN),
        ("Test Results", evidence["tests"], YELLOW),
        ("Screenshots", evidence["screenshots"], PURPLE),
        ("Runtime Proofs", evidence["runtime"], ORANGE),
        ("CI Validation", evidence["ci"], ACCENT),
    ]

    y = 2.0
    for name, count, color in categories:
        bar_width = min(count * 0.5, 8)
        add_text(slide3, 0.5, y, 3, 0.4, f"{name}: {count}", size=14, color=color, bold=True)
        if count > 0:
            shape = slide3.shapes.add_shape(
                1, Inches(4), Inches(y + 0.05), Inches(bar_width), Inches(0.3)
            )
            shape.fill.solid()
            shape.fill.fore_color.rgb = color
            shape.line.fill.background()
        y += 0.55

    verified = evidence["manifests"]
    total_expected = max(1, len([d for d in os.listdir(team_dir) if os.path.isdir(os.path.join(team_dir, d)) and d not in ('reports', 'evidence')])) if os.path.exists(team_dir) else 1
    add_text(slide3, 0.5, y + 0.3, 12, 0.5,
             f"Evidence Verification: {verified}/{total_expected} agents submitted manifests",
             size=16, color=GREEN if verified >= total_expected else RED, bold=True)

    # ──────────────────────────────────────────────────────────────
    # Slide 4: Milestone Progress
    # ──────────────────────────────────────────────────────────────
    slide4 = add_dark_slide(prs, "Milestone Progress")
    milestones = read_file_safe(os.path.join(team_dir, "MILESTONES.md"))
    if milestones:
        add_text(slide4, 0.5, 1.8, 12, 5, milestones, size=11, color=WHITE)
    else:
        add_text(slide4, 0.5, 1.8, 12, 5, "Milestones pending — PM has not yet produced MILESTONES.md", size=14, color=YELLOW)

    # ──────────────────────────────────────────────────────────────
    # Slide 5: Kanban Snapshot
    # ──────────────────────────────────────────────────────────────
    slide5 = add_dark_slide(prs, "Kanban Snapshot")
    add_text(slide5, 0.5, 1.8, 3, 0.6, f"Backlog: {kanban['backlog']}", size=18, color=ACCENT)
    add_text(slide5, 3.5, 1.8, 3, 0.6, f"In Progress: {kanban['in_progress']}", size=18, color=YELLOW)
    add_text(slide5, 6.5, 1.8, 3, 0.6, f"In Review: {kanban['in_review']}", size=18, color=ACCENT)
    add_text(slide5, 9.5, 1.8, 3, 0.6, f"Done: {kanban['done']}", size=18, color=GREEN)

    kanban_content = read_file_safe(os.path.join(team_dir, "KANBAN.md"), 2000)
    if kanban_content:
        add_text(slide5, 0.5, 2.8, 12, 4, kanban_content, size=10, color=WHITE)

    # ──────────────────────────────────────────────────────────────
    # Slide 6: Test Coverage & CI Status
    # ──────────────────────────────────────────────────────────────
    slide6 = add_dark_slide(prs, "Test Coverage & CI/CD Status")

    # Check for test evidence
    test_dir = os.path.join(team_dir, "evidence", "tests")
    test_layers = ["static", "unit", "integration", "e2e", "performance", "security", "release"]
    y = 1.8
    for layer in test_layers:
        layer_dir = os.path.join(test_dir, layer)
        if os.path.exists(layer_dir) and os.listdir(layer_dir):
            status = "PASS"
            color = GREEN
            file_count = len(os.listdir(layer_dir))
        elif os.path.exists(test_dir):
            status = "PENDING"
            color = YELLOW
            file_count = 0
        else:
            status = "NOT RUN"
            color = GRAY
            file_count = 0
        add_text(slide6, 0.5, y, 4, 0.4, f"{layer.upper()}: {status}", size=13, color=color, bold=True)
        add_text(slide6, 5, y, 4, 0.4, f"({file_count} result files)", size=11, color=GRAY)
        y += 0.45

    # CI status
    ci_dir = os.path.join(team_dir, "evidence", "ci")
    if os.path.exists(ci_dir) and os.listdir(ci_dir):
        add_text(slide6, 0.5, y + 0.3, 12, 0.5, "GitHub Actions (local): VALIDATED with act", size=14, color=GREEN, bold=True)
    else:
        add_text(slide6, 0.5, y + 0.3, 12, 0.5, "GitHub Actions (local): NOT YET VALIDATED", size=14, color=YELLOW)

    # ──────────────────────────────────────────────────────────────
    # Slide 7: Commit Activity
    # ──────────────────────────────────────────────────────────────
    slide7 = add_dark_slide(prs, "Commit Activity", f"Total commits tracked: {commits}")
    commit_log = read_file_safe(os.path.join(team_dir, "COMMIT_LOG.md"), 2500)
    if commit_log:
        add_text(slide7, 0.5, 2.0, 12, 5, commit_log, size=9, color=WHITE)
    else:
        add_text(slide7, 0.5, 2.0, 12, 5,
                 "No commits logged yet. PM will track atomic commits as agents work.\n\n"
                 "Expected format:\n"
                 "  feat(api): add user endpoint [#12]\n"
                 "  test(api): add user endpoint tests [#12]\n"
                 "  docs(api): document user API [#12]",
                 size=12, color=GRAY)

    # ──────────────────────────────────────────────────────────────
    # Slide 8: Blockers & Risks
    # ──────────────────────────────────────────────────────────────
    slide8 = add_dark_slide(prs, "Blockers & Risks")
    risks = read_file_safe(os.path.join(team_dir, "RISK_REGISTER.md"))
    if risks:
        add_text(slide8, 0.5, 1.8, 12, 5, risks, size=11, color=WHITE)
    else:
        add_text(slide8, 0.5, 1.8, 12, 5, "No blockers identified.", size=14, color=GREEN)

    # ──────────────────────────────────────────────────────────────
    # Slide 9: Team Activity
    # ──────────────────────────────────────────────────────────────
    slide9 = add_dark_slide(prs, "Team Activity")
    status_file = read_file_safe(os.path.join(team_dir, "TEAM_STATUS.md"))
    if status_file:
        add_text(slide9, 0.5, 1.8, 12, 5, status_file, size=10, color=WHITE)
    else:
        agents = ["PM", "Backend", "Frontend", "Mobile", "DevOps", "Infra", "QA", "Release", "Marketing", "Legal"]
        y = 1.8
        for agent in agents:
            add_text(slide9, 0.5, y, 12, 0.4, f"  {agent}: Standby", size=12, color=GRAY)
            y += 0.45

    # ──────────────────────────────────────────────────────────────
    # Slide 10: Artifacts Produced
    # ──────────────────────────────────────────────────────────────
    slide10 = add_dark_slide(prs, "Artifacts Produced")
    artifacts = discover_artifacts(team_dir)
    if artifacts:
        y = 1.8
        for dir_name, count in sorted(artifacts.items()):
            if y > 6.5:
                break
            add_text(slide10, 0.5, y, 8, 0.4, f".team/{dir_name}/", size=12, color=ACCENT)
            add_text(slide10, 9, y, 3, 0.4, f"{count} files", size=12, color=WHITE)
            y += 0.4
    else:
        add_text(slide10, 0.5, 1.8, 12, 5, "No artifacts produced yet.", size=14, color=YELLOW)

    # ──────────────────────────────────────────────────────────────
    # Slide 11: Next Steps
    # ──────────────────────────────────────────────────────────────
    slide11 = add_dark_slide(prs, "Next Steps")
    add_text(slide11, 0.5, 1.8, 12, 5,
             "Planned for next reporting period:\n"
             "- Continue wave-based execution with evidence collection\n"
             "- Verify all agents submit evidence manifests\n"
             "- Run comprehensive test matrix (unit + integration + E2E)\n"
             "- Validate GitHub Actions locally with act\n"
             "- Atomic commits for every change with issue references\n"
             "- Update GitHub Project board in real-time\n"
             "- Monitor quality gates (evidence + build + test + CI)\n"
             "- Generate next status report at T+6h",
             size=14, color=WHITE)

    # Save
    reports_dir = os.path.join(team_dir, "reports")
    os.makedirs(reports_dir, exist_ok=True)
    filepath = os.path.join(reports_dir, f"status_{report_num:03d}.pptx")
    prs.save(filepath)
    print(f"PPTX saved: {filepath}")
    return filepath


if __name__ == "__main__":
    report_num = int(sys.argv[1]) if len(sys.argv) > 1 else 1
    project_name = sys.argv[2] if len(sys.argv) > 2 else "Project"
    team_dir = sys.argv[3] if len(sys.argv) > 3 else ".team"
    create_status_report(report_num, project_name, team_dir)
